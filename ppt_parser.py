import os
from dataclasses import dataclass
from typing import List
from pathlib import Path
import subprocess
import tempfile


@dataclass
class PPTSlide:
    slide_number: int
    image_path: str  # Path to exported image
    notes_text: str  # Notes text content


def parse_ppt(pptx_path: str, output_dir: str) -> List[PPTSlide]:
    """
    Parse PPTX file and extract slides with notes.
    """
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")
    
    os.makedirs(output_dir, exist_ok=True)
    
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ImportError(
            "Missing required dependencies. Please install python-pptx"
        ) from e
    
    prs = Presentation(pptx_path)
    slides = []
    
    # 第一步：先提取所有备注文本
    for idx, slide in enumerate(prs.slides, start=1):
        notes_text = ""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text:
                notes_text = notes_frame.text.strip()
        slides.append(PPTSlide(
            slide_number=idx,
            image_path="",
            notes_text=notes_text
        ))
    
    # 第二步：尝试导出幻灯片图片
    # 方法 1：使用 PowerPoint COM（Windows）
    image_files = _try_export_with_powerpoint(pptx_path, output_dir, len(slides))
    
    # 方法 2：如果没有成功，尝试使用 LibreOffice
    if not image_files:
        image_files = _try_export_with_libreoffice(pptx_path, output_dir, len(slides))
    
    # 方法 3：如果都没有成功，创建简单占位图片
    if not image_files:
        from PIL import Image
        for idx in range(len(slides)):
            image_path = os.path.join(output_dir, f"slide_{idx+1:03d}.png")
            img = Image.new('RGB', (1280, 720), color='white')
            img.save(image_path)
            image_files.append(image_path)
    
    # 更新图片路径
    for idx in range(len(slides)):
        slides[idx].image_path = image_files[idx]
    
    return slides


def _try_export_with_powerpoint(pptx_path: str, output_dir: str, num_slides: int) -> List[str]:
    """
    使用 Windows PowerPoint COM 接口导出幻灯片为图片
    """
    image_files = []
    
    if os.name != 'nt':
        return []
    
    try:
        import win32com.client as win32
        import pythoncom
        
        pythoncom.CoInitialize()
        
        print("正在使用 PowerPoint 导出幻灯片...")
        
        powerpoint = None
        presentation = None
        
        try:
            # 创建 PowerPoint 应用实例
            powerpoint = win32.Dispatch("PowerPoint.Application")
            powerpoint.Visible = True  # 必须可见才能运行某些操作
            
            # 打开演示文稿
            full_path = os.path.abspath(pptx_path)
            presentation = powerpoint.Presentations.Open(full_path, WithWindow=False)
            
            # 导出每张幻灯片
            for slide_idx in range(1, presentation.Slides.Count + 1):
                slide = presentation.Slides(slide_idx)
                output_path = os.path.join(output_dir, f"slide_{slide_idx:03d}.png")
                
                # Export 参数: (FileName, FilterName, ScaleWidth, ScaleHeight)
                slide.Export(output_path, "PNG", 1920, 1080)
                
                if os.path.exists(output_path):
                    image_files.append(output_path)
            
            print(f"成功导出 {len(image_files)} 张幻灯片图片")
            
        finally:
            if presentation:
                presentation.Close()
            if powerpoint:
                powerpoint.Quit()
            pythoncom.CoUninitialize()
        
        if len(image_files) == num_slides:
            return image_files
        
        return []
        
    except Exception as e:
        print(f"PowerPoint 导出失败: {e}")
        return []


def _try_export_with_libreoffice(pptx_path: str, output_dir: str, num_slides: int) -> List[str]:
    """
    尝试使用 LibreOffice 导出 PPTX 为图片
    """
    image_files = []
    try:
        soffice_path = None
        # Windows 常见路径
        if os.name == 'nt':
            common_paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
            ]
            for path in common_paths:
                if os.path.exists(path):
                    soffice_path = path
                    break
        
        if not soffice_path:
            soffice_path = "soffice"
        
        # 使用 LibreOffice 转换为 PNG
        output_format = "png"
        cmd = [
            soffice_path,
            "--headless",
            "--convert-to", output_format,
            "--outdir", output_dir,
            pptx_path
        ]
        
        subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        
        # 检查生成的图片
        base_name = Path(pptx_path).stem
        for idx in range(1, num_slides + 1):
            # 尝试不同的命名格式
            possible_names = [
                f"{base_name}.png",
                f"{base_name}{idx}.png",
                f"{base_name}{idx-1}.png",
            ]
            for name in possible_names:
                possible_path = os.path.join(output_dir, name)
                if os.path.exists(possible_path):
                    new_path = os.path.join(output_dir, f"slide_{idx:03d}.png")
                    os.rename(possible_path, new_path)
                    image_files.append(new_path)
                    break
    
    except Exception as e:
        print(f"LibreOffice 导出失败: {e}")
    
    if len(image_files) == num_slides:
        return image_files
    
    return []
